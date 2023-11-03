from pptx import Presentation
from pptx.util import Inches
import os
import re

# 创建一个演示文稿对象
prs = Presentation()

# 设置幻灯片宽高比为16:9
prs.slide_width = Inches(10)  # PowerPoint的宽高比16:9的默认宽度
prs.slide_height = Inches(10 * 9.0 / 16.0)  # 根据16:9的比例计算高度

# 指定包含PNG图片的目录
image_dir = 'video'

# 定义一个提取数字的函数
def extract_number(filename):
    numbers = re.findall(r'\d+', filename)
    return int(numbers[0]) if numbers else None

# 获取目录中所有png文件，并按数字排序
images = sorted(
    (img for img in os.listdir(image_dir) if img.endswith('.png')),
    key=extract_number
)

# 遍历排序后的图片文件
for image_filename in images:
    # 添加一张新幻灯片
    slide_layout = prs.slide_layouts[6]  # 使用空白幻灯片模板
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置图片路径
    img_path = os.path.join(image_dir, image_filename)
    
    # 在当前幻灯片中添加图片，图片填充整个幻灯片
    left = top = Inches(0)
    slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)

# 保存演示文稿
prs.save('output_presentation.pptx')
